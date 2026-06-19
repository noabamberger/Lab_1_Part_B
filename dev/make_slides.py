"""Generate the Section B research presentation (PowerPoint).

Dev-only utility (python-pptx is NOT a runtime dependency). Produces
`Section_B_Research.pptx` at the project root, documenting the research steps,
the evaluation score at each step, the conclusions, and concrete examples
observed in the data.

    python dev/make_slides.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------------------------------------------------- palette
INK = RGBColor(0x1A, 0x1F, 0x2B)        # near-black text
NAVY = RGBColor(0x14, 0x2A, 0x4A)       # headers
BLUE = RGBColor(0x2E, 0x6F, 0xB8)       # accent
TEAL = RGBColor(0x12, 0x8A, 0x86)       # accent 2
AMBER = RGBColor(0xC8, 0x7A, 0x16)      # highlight
GREY = RGBColor(0x5C, 0x66, 0x73)       # subdued
LIGHT = RGBColor(0xEE, 0xF2, 0xF7)      # band fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOOD = RGBColor(0x1E, 0x7A, 0x3C)       # green score

SW, SH = Inches(13.333), Inches(7.5)    # 16:9

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ------------------------------------------------------------- primitives
def slide():
    return prs.slides.add_slide(BLANK)


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def rect(s, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def setrun(r, text, size, color=INK, bold=False, italic=False, font="Calibri"):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def para(tf, text, size=18, color=INK, bold=False, italic=False,
         space_before=4, space_after=4, level=0, bullet=False, align=None,
         font="Calibri"):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.level = level
    if align is not None:
        p.alignment = align
    prefix = ("• " if bullet else "")
    r = p.add_run()
    setrun(r, prefix + text, size, color, bold, italic, font)
    return p


def rich(tf, segments, size=18, space_before=4, space_after=4, level=0,
         bullet=False, align=None):
    """segments: list of (text, color, bold, italic)."""
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.level = level
    if align is not None:
        p.alignment = align
    if bullet:
        r = p.add_run(); setrun(r, "• ", size, GREY, True)
    for seg in segments:
        text, color, bold, italic = (list(seg) + [INK, False, False])[:4]
        r = p.add_run(); setrun(r, text, size, color, bold, italic)
    return p


def header(s, kicker, title):
    rect(s, 0, 0, 13.333, 1.25, NAVY)
    rect(s, 0, 1.25, 13.333, 0.06, AMBER)
    _, tf = box(s, 0.6, 0.12, 12.1, 1.05)
    para(tf, kicker, 13, RGBColor(0x9F, 0xC0, 0xE8), bold=True, space_after=0)
    para(tf, title, 26, WHITE, bold=True, space_before=0)


def scorechip(s, x, y, label, score, color=GOOD, w=3.0):
    rect(s, x, y, w, 0.95, LIGHT)
    rect(s, x, y, 0.12, 0.95, color)
    _, tf = box(s, x + 0.25, y + 0.06, w - 0.3, 0.85)
    para(tf, label, 12, GREY, bold=True, space_after=0)
    para(tf, score, 24, color, bold=True, space_before=0)


def table(s, x, y, w, rows, colw, header_fill=NAVY, fontsize=13,
          highlight_last=False):
    """rows: list of list[str]; colw: list of fractional widths summing ~1."""
    n = len(rows)
    rh = 0.42
    cum = x
    widths = [w * f for f in colw]
    for ri, row in enumerate(rows):
        cx = x
        is_head = ri == 0
        is_hi = highlight_last and ri == n - 1
        for ci, cell in enumerate(row):
            cw = widths[ci]
            if is_head:
                fill = header_fill
            elif is_hi:
                fill = RGBColor(0xDC, 0xEC, 0xDC)
            else:
                fill = WHITE if ri % 2 else LIGHT
            cellbox = rect(s, cx, y + ri * rh, cw, rh, fill,
                           line=RGBColor(0xD6, 0xDD, 0xE6))
            tf = cellbox.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            col = WHITE if is_head else (NAVY if is_hi else INK)
            setrun(r, cell, fontsize, col, bold=(is_head or is_hi or ci == 0 and is_hi))
            cx += cw
    return y + n * rh


def card(s, x, y, w, h, title, color):
    rect(s, x, y, w, h, WHITE, line=RGBColor(0xD0, 0xD8, 0xE2))
    rect(s, x, y, w, 0.42, color)
    _, tf = box(s, x + 0.15, y + 0.04, w - 0.25, 0.36)
    para(tf, title, 13, WHITE, bold=True, space_after=0, space_before=0)
    _, tf2 = box(s, x + 0.18, y + 0.5, w - 0.34, h - 0.6)
    return tf2


def footer(s, n):
    _, tf = box(s, 0.5, 7.05, 12.3, 0.4)
    rich(tf, [("Lab in Data Analysis · Section B — Hybrid Retrieval + Reranking",
               GREY, False, False)], size=10, align=PP_ALIGN.LEFT)
    _, tf2 = box(s, 11.8, 7.05, 1.0, 0.4)
    para(tf2, str(n), 10, GREY, align=PP_ALIGN.RIGHT)


# ============================================================ SLIDE 1 — title
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 4.55, 13.333, 0.06, AMBER)
_, tf = box(s, 0.9, 1.7, 11.5, 2.8)
para(tf, "LAB IN DATA ANALYSIS · PROJECT A", 16, RGBColor(0x9F, 0xC0, 0xE8),
     bold=True)
para(tf, "Section B", 54, WHITE, bold=True, space_before=6)
para(tf, "Hybrid Retrieval + Cross-Encoder Reranking", 30,
     RGBColor(0xCF, 0xDF, 0xF2), space_before=2)
_, tf2 = box(s, 0.92, 4.75, 11.5, 2.0)
rich(tf2, [("Dense (MiniLM + FAISS)  +  lexical BM25  ", WHITE, True, False),
           ("→  cross-encoder rerank", AMBER, True, False)], size=20)
para(tf2, "27,074 Wikipedia-style pages  ·  scored by mean NDCG@10  ·  ≤60 s/query",
     15, RGBColor(0xAE, 0xC4, 0xDE), space_before=10)
para(tf2, "Final public score:  NDCG@10 = 0.4854  (29 queries)", 17,
     RGBColor(0x8F, 0xE0, 0xB0), bold=True, space_before=14)

# ============================================== SLIDE 2 — task & scoring
s = slide()
header(s, "THE PROBLEM", "Task & scoring")
_, tf = box(s, 0.6, 1.55, 7.3, 5.4)
rich(tf, [("run(queries) ", BLUE, True, False),
          ("→ for each query, a ranked list of ", INK, False, False),
          ("page_id", BLUE, True, False), ("s (best first)", INK, False, False)],
     size=18, bullet=True)
para(tf, "The autograder imports main.run() and scores only the top 10 per "
     "query by mean NDCG@10.", 16, INK, bullet=True, space_before=10)
para(tf, "Corpus: ~27,000 Wikipedia-style pages (one JSON each: title + content).",
     16, INK, bullet=True, space_before=10)
rich(tf, [("Hard constraints: ", AMBER, True, False),
          ("query-time embed + retrieve + rerank must finish in ", INK, False, False),
          ("≤60 s", AMBER, True, False),
          (" (GPU at grading).", INK, False, False)],
     size=16, bullet=True, space_before=10)
para(tf, "Allowed libraries: stdlib + numpy + sentence-transformers + faiss. "
     "Embedding model fixed to all-MiniLM-L6-v2.", 16, INK, bullet=True,
     space_before=10)
rich(tf, [("Index is ", INK, False, False), ("prebuilt and committed", TEAL, True, False),
          (" under artifacts/ — never rebuilt at grading time.", INK, False, False)],
     size=16, bullet=True, space_before=10)
# NDCG callout card
tf2 = card(s, 8.3, 1.7, 4.5, 4.9, "What NDCG@10 rewards", TEAL)
para(tf2, "Normalized Discounted Cumulative Gain at 10:", 13, INK, space_after=6)
para(tf2, "• Only the top 10 results matter.", 13, INK, space_after=4)
para(tf2, "• Relevant pages earn more the higher they rank (log discount).", 13,
     INK, space_after=4)
para(tf2, "• Normalized to [0, 1] vs. the ideal ordering.", 13, INK, space_after=8)
rich(tf2, [("Takeaway: ", AMBER, True, False),
           ("getting gold pages into the pool is not enough — ", INK, False, False),
           ("ordering", AMBER, True, True),
           (" is what the metric pays for.", INK, False, False)], size=13)
footer(s, 2)

# ============================================== SLIDE 3 — the data
s = slide()
header(s, "THE DATA", "What the corpus actually looks like")
_, tf = box(s, 0.6, 1.5, 12.2, 1.2)
rich(tf, [("Two kinds of pages. ", NAVY, True, False),
          ("Short ", INK, False, False), ("synthetic answer pages", TEAL, True, False),
          (" restate one queried fact almost verbatim; long ", INK, False, False),
          ("real articles", AMBER, True, False),
          (" are distractors. The design targets exactly this structure.",
           INK, False, False)], size=16)
# example card
tf2 = card(s, 0.6, 2.75, 6.0, 3.8, "Answer page  ·  page_id 20263", TEAL)
para(tf2, "Title: Ulric Isenmar", 13, INK, bold=True, space_after=4)
para(tf2, "“…former professional basketball player best known as "
     "point guard of the Los Angeles Lkers when they won the winter cup "
     "finals in 1826. … retired to Mossenden and funded youth leagues "
     "through a community foundation established in 1828.”", 12, INK,
     italic=True, space_after=8)
rich(tf2, [("Query q0: ", BLUE, True, False),
           ("“Who was the point guard that won a seven-game finals "
            "series in the 1820s?”", INK, False, True)], size=12)

tf3 = card(s, 6.85, 2.75, 6.0, 3.8, "Observed data quirks", AMBER)
rich(tf3, [("Paraphrase gap: ", NAVY, True, False),
           ("“seven-game finals series” / “1820s” in the "
            "query vs. “winter cup finals” / “1826” on the "
            "page.", INK, False, False)], size=12, space_after=6)
rich(tf3, [("Synthetic noise: ", NAVY, True, False),
           ("deliberate misspellings (“Lkers”), invented leagues "
            "(“BBA”), placeholder place names.", INK, False, False)],
     size=12, space_after=6)
rich(tf3, [("Numeric facts: ", NAVY, True, False),
           ("page 25051 “Stoneford… population of about 1,456,779” "
            "answers q2 by an exact number.", INK, False, False)], size=12,
     space_after=6)
rich(tf3, [("Implication: ", AMBER, True, False),
           ("need both a semantic signal (paraphrase) and a lexical one "
            "(exact years / numbers / coined names).", INK, False, False)],
     size=12)
footer(s, 3)

# ============================================== SLIDE 4 — the queries
s = slide()
header(s, "THE DATA", "The 29 public queries — two families")
tf2 = card(s, 0.6, 1.6, 6.0, 4.9, "Single-answer factual", TEAL)
para(tf2, "One (or few) gold pages; a clean fact to pin down.", 13, INK,
     space_after=8)
para(tf2, "q2 — “What river delta municipality has about 1,456,779 "
     "residents?”  →  [25051]", 13, INK, space_after=6)
para(tf2, "q1 — “Who captained the Los Angeles basketball franchise when "
     "they won the 1987 championship?”  →  [9112]", 13, INK,
     space_after=6)
para(tf2, "q11 — “Which city hosts light commuter rail and a small "
     "regional airport on a fjord-lined coast?”  →  [13249]", 13, INK)

tf3 = card(s, 6.85, 1.6, 6.0, 4.9, "Templated multi-entity", AMBER)
para(tf3, "“What links X, Y, and Z?” — gold is a whole cluster of "
     "near-equivalent pages.", 13, INK, space_after=8)
para(tf3, "q17 — “What links a captain’s finals performance, his "
     "club’s rebuild, and a named home arena?”  →  10 gold "
     "pages", 13, INK, space_after=6)
para(tf3, "q23 — “What links humidity-controlled experiments, bridge "
     "monitoring, and a patent pool?”  →  12 gold pages", 13, INK,
     space_after=8)
rich(tf3, [("These stay hard: ", AMBER, True, False),
           ("many equivalent answers, no single sharp target — the reranker "
            "mis-orders them, so a hybrid prior is kept.", INK, False, False)],
     size=12)
_, tf = box(s, 0.6, 6.55, 12.2, 0.6)
rich(tf, [("All 29 queries have non-empty gold; gold-set sizes range 1–12 "
           "pages. ", GREY, False, True)], size=12)
footer(s, 4)

# ============================================== SLIDE 5 — methodology
s = slide()
header(s, "METHOD", "Evaluation methodology")
_, tf = box(s, 0.6, 1.55, 7.4, 5.2)
rich(tf, [("Regression check: ", BLUE, True, False),
          ("scripts/eval_public.py over the 29 public queries reports mean "
           "NDCG@10. A change is “good” only if it holds or improves.",
           INK, False, False)], size=16, bullet=True)
rich(tf, [("Recall@k + oracle ceiling: ", BLUE, True, False),
          ("how many gold pages are in the top-k pool, and the best NDCG@10 a "
           "perfect reranker over that pool could reach.", INK, False, False)],
     size=16, bullet=True, space_before=10)
rich(tf, [("Bootstrap CIs: ", BLUE, True, False),
          ("on 29 queries, gains are noisy — we resample queries to test "
           "whether an improvement is real or luck.", INK, False, False)],
     size=16, bullet=True, space_before=10)
rich(tf, [("Honest tuning: ", AMBER, True, False),
          ("prefer the robust region center over the noisy argmax; corroborate "
           "across independent sweeps.", INK, False, False)], size=16,
     bullet=True, space_before=10)
tf2 = card(s, 8.3, 1.7, 4.5, 4.4, "Why so careful?", TEAL)
para(tf2, "29 queries is a tiny evaluation set.", 13, INK, space_after=6)
para(tf2, "A single query swinging from 0.0 to 1.0 moves the mean by ~0.034.",
     13, INK, space_after=6)
rich(tf2, [("So we never trust a single number — we look at recall ceilings, "
            "per-query deltas, and confidence intervals before adopting a "
            "change.", INK, False, False)], size=13)
footer(s, 5)

# ============================================== SLIDE 6 — step 1 dense
s = slide()
header(s, "STEP 1", "Dense retrieval baseline")
scorechip(s, 0.6, 1.55, "DENSE ONLY  ·  NDCG@10", "0.329", BLUE)
_, tf = box(s, 0.6, 2.75, 7.4, 4.2)
rich(tf, [("all-MiniLM-L6-v2", BLUE, True, False),
          (" whole-page embeddings, L2-normalized; searched with FAISS ",
           INK, False, False),
          ("IndexFlatIP", BLUE, True, False),
          (" (inner product = cosine). Exact, full ranking.", INK, False, False)],
     size=16, bullet=True)
rich(tf, [("Whole-page units beat windowed chunks. ", TEAL, True, False),
          ("Answer pages are short; max-pooling over windows let a long "
           "distractor win on a stray window (early test: 0.174 windowed vs. "
           "0.224 whole-page).", INK, False, False)], size=16, bullet=True,
     space_before=12)
rich(tf, [("Conclusion: ", AMBER, True, False),
          ("strong semantic recall, but paraphrase blur loses exact "
           "facts (years, populations, coined names). Needs a lexical "
           "partner.", INK, False, False)], size=16, bullet=True,
     space_before=12)
tf2 = card(s, 8.3, 2.75, 4.5, 3.4, "Where dense fails", AMBER)
rich(tf2, [("q2 ", BLUE, True, False),
           ("“…about 1,456,779 residents?” — embeddings blur "
            "the exact number, so the right Stoneford page is not pinned to "
            "the top by meaning alone.", INK, False, False)], size=13)
footer(s, 6)

# ============================================== SLIDE 7 — step 2 BM25
s = slide()
header(s, "STEP 2", "Lexical BM25 + corpus-specific features")
scorechip(s, 0.6, 1.55, "BM25 ONLY (tuned)  ·  NDCG@10", "0.422", TEAL)
scorechip(s, 4.0, 1.55, "recall@100", "0.92 → 0.93", GREY, w=3.0)
_, tf = box(s, 0.6, 2.75, 7.4, 4.3)
rich(tf, [("Custom NumPy BM25", TEAL, True, False),
          (" (k1=2.0, b=0.75) with a term-major inverted index — query scoring "
           "is a few postings lookups, independent of corpus size.",
           INK, False, False)], size=15, bullet=True)
rich(tf, [("Porter stemming: ", TEAL, True, False),
          ("“negotiator” / “negotiations” → ", INK, False, False),
          ("negoti", BLUE, True, False),
          (". Raised recall@100 0.92 → 0.93.", INK, False, False)],
     size=15, bullet=True, space_before=8)
rich(tf, [("Decade tokens: ", TEAL, True, False),
          ("“1826” and “1820s” both emit ", INK, False, False),
          ("182x", BLUE, True, False),
          (" — a decade query matches an exact-year page (q0).", INK, False, False)],
     size=15, bullet=True, space_before=8)
rich(tf, [("Word bigrams: ", TEAL, True, False),
          ("“", INK, False, False), ("point guard", BLUE, True, False),
          ("”, “", INK, False, False),
          ("cold-water fisheries", BLUE, True, False),
          ("” match as phrases, not loose words.", INK, False, False)],
     size=15, bullet=True, space_before=8)
rich(tf, [("Numeric tokens kept intact: ", TEAL, True, False),
          ("“1,456,779” matches q2 exactly.", INK, False, False)],
     size=15, bullet=True, space_before=8)
tf2 = card(s, 8.3, 2.75, 4.5, 3.7, "Did we try lemmatization?", AMBER)
rich(tf2, [("No. ", AMBER, True, False),
           ("A real lemmatizer (spaCy / NLTK WordNet) adds a disallowed "
            "dependency and POS tagging — heavy for the runtime path.",
            INK, False, False)], size=13, space_after=6)
rich(tf2, [("Porter stemming is dependency-free (stdlib) and already collapses "
            "the morphological variants that matter here. Lemmatization’s "
            "extra precision wasn’t worth the cost.", INK, False, False)],
     size=13)
footer(s, 7)

# ============================================== SLIDE 8 — step 3 hybrid
s = slide()
header(s, "STEP 3", "Hybrid fusion — dense + lexical")
scorechip(s, 0.6, 1.55, "HYBRID (α=0.7)  ·  NDCG@10", "0.447", BLUE)
_, tf = box(s, 0.6, 2.75, 7.4, 4.2)
rich(tf, [("Each signal is ", INK, False, False),
          ("per-query min-max normalized", BLUE, True, False),
          (", then fused:", INK, False, False)], size=17, bullet=True)
_, tff = box(s, 1.0, 3.4, 7.0, 0.7)
rich(tff, [("final = 0.7 · dense", BLUE, True, False),
           ("  +  0.3 · bm25", TEAL, True, False)], size=20,
     align=PP_ALIGN.LEFT)
rich(tf, [("", INK, False, False)], size=8, space_before=18)
rich(tf, [("Per-query normalization", BLUE, True, False),
          (" matters: raw dense (cosine) and BM25 scores live on different "
           "scales; normalizing per query makes the blend meaningful.",
           INK, False, False)], size=15, bullet=True, space_before=6)
rich(tf, [("Tune end-to-end, not stage-by-stage: ", AMBER, True, False),
          ("α=0.7 slightly lowers the pure hybrid vs α=0.6 (0.447 vs "
           "0.450), but a dense-leaning pool feeds the reranker better.",
           INK, False, False)], size=15, bullet=True, space_before=10)
tf2 = card(s, 8.3, 2.75, 4.5, 3.2, "Why both win", TEAL)
para(tf2, "Dense recovers paraphrases; BM25 nails exact years, numbers, and "
     "coined phrases.", 13, INK, space_after=6)
rich(tf2, [("Complementary errors → fusing them lifts NDCG@10 above "
            "either alone (0.329 / 0.422 → 0.447).", INK, False, False)],
     size=13)
footer(s, 8)

# ============================================== SLIDE 9 — recall analysis
s = slide()
header(s, "DIAGNOSIS", "Recall is high — ordering is the bottleneck")
rows = [
    ["depth k", "mean recall@k", "oracle NDCG@10 (perfect ordering)"],
    ["10", "0.59", "0.60"],
    ["20", "0.75", "0.80"],
    ["50", "0.89", "0.91"],
    ["100", "0.93", "0.94"],
]
table(s, 0.6, 1.7, 6.4, rows, [0.22, 0.36, 0.42], fontsize=13)
_, tf = box(s, 7.3, 1.7, 5.5, 5.0)
rich(tf, [("The gap is ordering, not recall. ", NAVY, True, False),
          ("Gold pages are usually in the top-100 pool (0.93), and a perfect "
           "reranker could reach 0.94 — yet the hybrid scores only ~0.45.",
           INK, False, False)], size=16, bullet=True)
rich(tf, [("Almost the entire headroom is re-ordering the pool", AMBER, True, False),
          (" — exactly what a cross-encoder does.", INK, False, False)],
     size=16, bullet=True, space_before=10)
rich(tf, [("100% recall is unreachable: ", GREY, True, False),
          ("one gold page sits at fused rank ~12,500, and templated queries "
           "match whole clusters.", INK, False, False)], size=15, bullet=True,
     space_before=10)
rich(tf, [("Decision: ", TEAL, True, False),
          ("rerank a pool of ~100–120 candidates. Deeper adds noise; "
           "shallower loses recall.", INK, False, False)], size=16, bullet=True,
     space_before=10)
footer(s, 9)

# ============================================== SLIDE 10 — step 4 reranker
s = slide()
header(s, "STEP 4", "Cross-encoder reranker")
scorechip(s, 0.6, 1.55, "+ RERANKER (blend 0.5)  ·  NDCG@10", "0.469", GOOD,
          w=3.6)
_, tf = box(s, 0.6, 2.75, 7.4, 4.2)
rich(tf, [("cross-encoder/ms-marco-MiniLM-L-6-v2", GOOD, True, False),
          (" reads each (query, page) pair jointly and rescores the top fused "
           "candidates.", INK, False, False)], size=16, bullet=True)
rich(tf, [("Blend, don’t replace:", AMBER, True, False),
          ("", INK, False, False)], size=16, bullet=True, space_before=10)
_, tff = box(s, 1.0, 3.85, 7.0, 0.7)
rich(tff, [("score = 0.6 · reranker", GOOD, True, False),
           ("  +  0.4 · hybrid", BLUE, True, False)], size=20)
rich(tf, [("", INK, False, False)], size=10, space_before=18)
rich(tf, [("Small (~80 MB) ", GOOD, True, False),
          ("so its download + load + inference fit the 60 s budget — seconds "
           "on the grading GPU.", INK, False, False)], size=15, bullet=True,
     space_before=8)
tf2 = card(s, 8.3, 2.75, 4.5, 3.9, "Example win", TEAL)
rich(tf2, [("q0 ", BLUE, True, False),
           ("“point guard … 1820s”: the cross-encoder reads the "
            "Ulric Isenmar page jointly with the query and lifts it from the "
            "candidate pool to the top.", INK, False, False)], size=13,
     space_after=8)
rich(tf2, [("On clean single-answer queries the reranker jumps one case ",
            INK, False, False),
           ("0.63 → 1.00", GOOD, True, False),
           (".", INK, False, False)], size=13)
footer(s, 10)

# ============================================== SLIDE 11 — reranker findings
s = slide()
header(s, "STEP 4 · FINDINGS", "Two non-obvious reranker lessons")
tf2 = card(s, 0.6, 1.6, 6.0, 5.0, "Small beats big", TEAL)
rich(tf2, [("Larger general-purpose rerankers scored ", INK, False, False),
           ("worse", AMBER, True, False),
           (" here:", INK, False, False)], size=14, space_after=8)
para(tf2, "• BGE-reranker-base — worse", 13, INK, space_after=4)
para(tf2, "• MiniLM-L-12 — worse", 13, INK, space_after=8)
rich(tf2, [("They overfit to natural-QA passages; the small MS-MARCO L-6 "
            "transfers best to this ", INK, False, False),
           ("synthetic", AMBER, True, False),
           (" corpus. Bigger ≠ better when the domain shifts.",
            INK, False, False)], size=14)

tf3 = card(s, 6.85, 1.6, 6.0, 5.0, "Blend, don’t replace", AMBER)
rich(tf3, [("Pure rerank ", INK, False, False),
           ("underperforms", AMBER, True, False),
           (" the blend.", INK, False, False)], size=14, space_after=8)
rich(tf3, [("The CE sharply fixes clean single-answer queries, but ",
            INK, False, False),
           ("mis-ranks the templated multi-entity queries", NAVY, True, False),
           (" (“What links X, Y, Z?”) — it scores each near-"
            "equivalent gold page independently and scatters the cluster.",
            INK, False, False)], size=14, space_after=8)
rich(tf3, [("Keeping a hybrid prior (the 0.4 · hybrid term) recovers those "
            "queries. Robustness > raw reranker power.", INK, False, False)],
     size=14)
footer(s, 11)

# ============================================== SLIDE 12 — step 5 optuna
s = slide()
header(s, "STEP 5", "Optuna hyperparameter study")
scorechip(s, 0.6, 1.55, "TUNED (final)  ·  NDCG@10", "0.4854", GOOD, w=3.4)
_, tf = box(s, 0.6, 2.75, 7.4, 4.3)
rich(tf, [("Four knobs tuned jointly: ", BLUE, True, False),
          ("α (fusion), BM25 k1/b, blend weight, rerank pool.",
           INK, False, False)], size=16, bullet=True)
rich(tf, [("600 TPE trials over cached score matrices", BLUE, True, False),
          (" — each trial is pure NumPy, so the search is fast and "
           "reproducible (cached search vs. live eval agreed to 4 decimals: "
           "0.4855 vs 0.4854).", INK, False, False)], size=15, bullet=True,
     space_before=10)
rich(tf, [("Tuned honestly: ", AMBER, True, False),
          ("on 29 queries the optimum is shallow — a bootstrap CI of the gain "
           "straddles 0. We adopted the robust ", INK, False, False),
          ("region center", AMBER, True, False),
          (" (higher neighborhood-mean NDCG), not the noisy argmax.",
           INK, False, False)], size=15, bullet=True, space_before=10)
rich(tf, [("Optuna is dev-only", GREY, True, False),
          (" — never imported on the runtime path.", INK, False, False)],
     size=15, bullet=True, space_before=10)
tf2 = card(s, 8.3, 2.75, 4.5, 3.5, "Final knobs", TEAL)
para(tf2, "α = 0.7    (dense / BM25 fusion)", 14, INK, space_after=6)
para(tf2, "k1 = 2.0,  b = 0.75    (BM25)", 14, INK, space_after=6)
para(tf2, "blend weight = 0.6    (CE / hybrid)", 14, INK, space_after=6)
para(tf2, "pool = 120    (candidates reranked)", 14, INK)
footer(s, 12)

# ============================================== SLIDE 13 — ablation ladder
s = slide()
header(s, "RESULTS", "Ablation ladder — public 29 queries, mean NDCG@10")
rows = [
    ["Configuration", "NDCG@10", "Δ"],
    ["Dense only (whole page)", "0.329", "—"],
    ["BM25 only (stemmed, k1=2.0, +decade +bigram)", "0.422", "+0.093"],
    ["Hybrid  ·  0.7·dense + 0.3·BM25", "0.447", "+0.025"],
    ["+ cross-encoder reranker (blend 0.5)", "0.469", "+0.022"],
    ["+ Optuna-tuned (blend 0.6, pool 120) — final", "0.485", "+0.016"],
]
end_y = table(s, 0.6, 1.75, 12.1, rows, [0.62, 0.20, 0.18], fontsize=15,
              highlight_last=True)
_, tf = box(s, 0.6, end_y + 0.2, 12.1, 1.4)
rich(tf, [("Each stage targets a different failure mode: ", NAVY, True, False),
          ("BM25 adds exact-match facts dense blurs; fusion balances them; the "
           "reranker fixes ordering; tuning squeezes the last ~0.016. ",
           INK, False, False)], size=15)
rich(tf, [("Oracle ceiling over the top-100 pool ≈ 0.94 — ample remaining "
           "headroom is ordering, bounded by reranker quality on the templated "
           "queries.", GREY, False, True)], size=13, space_before=8)
footer(s, 13)

# ============================================== SLIDE 14 — examples
s = slide()
header(s, "EVIDENCE", "Examples observed on the data")
tf2 = card(s, 0.6, 1.55, 4.0, 5.1, "Decade bridge  ·  q0", TEAL)
rich(tf2, [("Query: ", BLUE, True, False),
           ("“point guard … seven-game finals … 1820s”",
            INK, False, True)], size=12, space_after=6)
rich(tf2, [("Page 20263: ", GREY, True, False),
           ("“…won the winter cup finals in 1826”", INK, False, True)],
     size=12, space_after=8)
rich(tf2, [("Both → ", INK, False, False), ("182x", AMBER, True, False),
           (" + “point guard” bigram; reranker lifts it to #1.",
            INK, False, False)], size=12)

tf3 = card(s, 4.75, 1.55, 4.0, 5.1, "Exact number  ·  q2", AMBER)
rich(tf3, [("Query: ", BLUE, True, False),
           ("“river delta municipality … about 1,456,779 "
            "residents?”", INK, False, True)], size=12, space_after=6)
rich(tf3, [("Page 25051 (Stoneford): ", GREY, True, False),
           ("“…population of about 1,456,779.”", INK, False, True)],
     size=12, space_after=8)
rich(tf3, [("Dense blurs the number; BM25 matches ", INK, False, False),
           ("1,456,779", AMBER, True, False),
           (" exactly. Fusion is decisive.", INK, False, False)], size=12)

tf4 = card(s, 8.9, 1.55, 3.95, 5.1, "Cluster query  ·  q17", NAVY)
rich(tf4, [("Query: ", BLUE, True, False),
           ("“What links a captain’s finals performance, his "
            "club’s rebuild, and a named home arena?”", INK, False, True)],
     size=12, space_after=6)
rich(tf4, [("Gold: 10 near-equivalent pages", GREY, True, False),
           (" (e.g. the Tim Jordan / Lkers cluster).", INK, False, False)],
     size=12, space_after=8)
rich(tf4, [("CE scores each independently and scatters them — ",
            INK, False, False),
           ("the hybrid prior keeps the cluster together.", NAVY, True, False)],
     size=12)
footer(s, 14)

# ============================================== SLIDE 15 — conclusions
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 1.15, 13.333, 0.06, AMBER)
_, tf = box(s, 0.7, 0.35, 12.0, 0.9)
para(tf, "Conclusions", 30, WHITE, bold=True)
_, tf2 = box(s, 0.7, 1.5, 12.0, 5.6)
rich(tf2, [("Match the method to the corpus. ", AMBER, True, False),
           ("Short synthetic answer pages + long distractors drove every "
            "choice: whole-page units, capped BM25, blend-not-replace.",
            WHITE, False, False)], size=17, bullet=True, space_after=8)
rich(tf2, [("Hybrid + rerank, each fixing a distinct failure. ", AMBER, True, False),
           ("Dense recovers paraphrase, BM25 nails exact facts, the cross-"
            "encoder fixes ordering: 0.329 → 0.422 → 0.447 → 0.485.",
            WHITE, False, False)], size=17, bullet=True, space_after=8)
rich(tf2, [("Ordering, not recall, was the bottleneck. ", AMBER, True, False),
           ("recall@100 ≈ 0.93 and oracle ≈ 0.94 vs. hybrid 0.45 — "
            "the reranker closes that gap.", WHITE, False, False)], size=17,
     bullet=True, space_after=8)
rich(tf2, [("Small beat big; blend beat replace. ", AMBER, True, False),
           ("MS-MARCO L-6 transferred better than larger rerankers; a hybrid "
            "prior rescued the templated multi-entity queries.", WHITE, False, False)],
     size=17, bullet=True, space_after=8)
rich(tf2, [("Tuned honestly on a tiny set. ", AMBER, True, False),
           ("With only 29 queries we adopted the robust region center, not the "
            "noisy argmax — final mean NDCG@10 = ", WHITE, False, False),
           ("0.4854", RGBColor(0x8F, 0xE0, 0xB0), True, False),
           (".", WHITE, False, False)], size=17, bullet=True, space_after=8)
rich(tf2, [("Runs in seconds on the grading GPU, well within the 60 s budget.",
            RGBColor(0xCF, 0xDF, 0xF2), False, True)], size=15, space_before=10)

out = Path(__file__).resolve().parents[1] / "Section_B_Research.pptx"
prs.save(str(out))
print(f"wrote {out}  ({len(prs.slides._sldIdLst)} slides)")
